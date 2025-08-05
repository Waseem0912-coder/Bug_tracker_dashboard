import pandas as pd
from pptx import Presentation
import PyPDF2
from PIL import Image
import pytesseract
import io

def process_files(files):
    processed = {}
    for f in files:
        filename = f.filename
        content = f.read()
        ext = filename.rsplit('.', 1)[-1].lower()
        if ext in ['xlsx', 'xls']:
            df = pd.read_excel(io.BytesIO(content))
            processed[filename] = {'type': 'excel', 'data': df}
        elif ext in ['pptx', 'ppt']:
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, 'text')])
                texts.append(slide_text)
            processed[filename] = {'type': 'ppt', 'data': texts}
        elif ext == 'pdf':
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = ' '.join([page.extract_text() for page in reader.pages])
            processed[filename] = {'type': 'pdf', 'data': text}
        elif ext in ['png', 'jpg', 'jpeg', 'bmp']:
            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
            processed[filename] = {'type': 'image', 'data': text}
        else:
            try:
                text = content.decode('utf-8', errors='ignore')
            except:
                text = str(content)
            processed[filename] = {'type': 'text', 'data': text}
    return processed

def process_text(text):
    # For now, just return the text
    return text
