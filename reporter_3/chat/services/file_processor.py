"""
This service handles the extraction of text content from various file types.
For text-based files (PDF, DOCX, etc.), it extracts the raw text.
For image files, it creates a simple text placeholder. The actual image data
will be encoded and sent to a vision model separately in the chat handler logic.
"""

import os
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from io import StringIO

# --- Individual File Processors ---

def process_pdf(file_path):
    """Extracts text from a PDF file."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            # Add a fallback for pages where text extraction might fail
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        filename = os.path.basename(file_path)
        return f"--- Content from PDF: {filename} ---\n{text}\n"
    except Exception as e:
        return f"Error processing PDF {os.path.basename(file_path)}: {e}"

def process_docx(file_path):
    """Extracts text from a Word (.docx) file."""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        filename = os.path.basename(file_path)
        return f"--- Content from DOCX: {filename} ---\n{text}\n"
    except Exception as e:
        return f"Error processing DOCX {os.path.basename(file_path)}: {e}"

def process_pptx(file_path):
    """Extracts text from a PowerPoint (.pptx) file."""
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        filename = os.path.basename(file_path)
        return f"--- Content from PPTX: {filename} ---\n{text}\n"
    except Exception as e:
        return f"Error processing PPTX {os.path.basename(file_path)}: {e}"

def process_excel(file_path):
    """Parses an Excel file and returns content as CSV-like strings for the LLM."""
    try:
        xls = pd.ExcelFile(file_path)
        filename = os.path.basename(file_path)
        full_content = f"--- Content from Excel: {filename} ---\n"
        for sheet_name in xls.sheet_names:
            # Read dataframe without filling NaNs to preserve original data state
            df = pd.read_excel(xls, sheet_name=sheet_name)
            # Convert dataframe to a CSV string, which is a great format for LLMs
            output = StringIO()
            df.to_csv(output, index=False)
            full_content += f"\n--- Sheet: {sheet_name} ---\n"
            full_content += output.getvalue()
        return full_content
    except Exception as e:
        return f"Error processing Excel file {os.path.basename(file_path)}: {e}"

def process_image_placeholder(file_path):
    """
    Creates a placeholder for an image. The actual image data will be handled
    separately and sent to a vision model like LLaVA.
    """
    filename = os.path.basename(file_path)
    return f"--- Image Uploaded: {filename} (Content will be processed by the selected vision model) ---\n"

def process_pasted_text(content):
    """Cleans up pasted text, specifically by stripping any HTML tags."""
    # Basic check to see if it's likely HTML to avoid unnecessary processing
    if '<' in content and '>' in content:
        soup = BeautifulSoup(content, 'html.parser')
        return soup.get_text(separator='\n', strip=True)
    return content

# --- Main Dispatcher and Aggregator ---

def process_uploaded_file(uploaded_file_obj):
    """
    Dispatcher function that takes an UploadedFile model instance,
    determines the file type, and calls the appropriate processor.
    Returns the textual representation or placeholder for the file.
    """
    file_path = uploaded_file_obj.file.path
    _, extension = os.path.splitext(file_path.lower())

    processors = {
        '.pdf': process_pdf,
        '.docx': process_docx,
        '.pptx': process_pptx,
        '.xlsx': process_excel,
        '.xls': process_excel,
        # Image types now point to the placeholder function
        '.png': process_image_placeholder,
        '.jpg': process_image_placeholder,
        '.jpeg': process_image_placeholder,
        '.tiff': process_image_placeholder,
        '.bmp': process_image_placeholder,
        '.gif': process_image_placeholder,
    }

    processor = processors.get(extension)
    if processor:
        return processor(file_path)
    else:
        # Fallback for plain text files or other unsupported types
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            filename = os.path.basename(file_path)
            return f"--- Content from Text File: {filename} ---\n{content}\n"
        except Exception as e:
            return f"Unsupported file type or error reading file {os.path.basename(file_path)}: {e}"

def aggregate_content(contents_list):
    """
    Combines content from multiple sources into a single string.
    This aggregated text will form the main context for the LLM.
    """
    header = "=== Aggregated Content for Analysis ===\n"
    # Filter out any None or empty string results before joining
    filtered_contents = [c for c in contents_list if c and c.strip()]
    return header + "\n".join(filtered_contents)